"""本地 PPT 生成服务

当讯飞 PPT API 不可用时的降级方案：
1. LLM 两阶段生成：大纲 → 逐页扩写详细内容
2. python-pptx 组装 .pptx 文件
"""

import json
import logging
import os
import re
import tempfile
import uuid

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from services.spark_service import spark_service

logger = logging.getLogger(__name__)

# ---- 字体：Windows 中文系统必备 ----
FONT_BODY = "微软雅黑"
FONT_TITLE = "微软雅黑"

# ---- 色彩体系 ----
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHITE  = RGBColor(0xF7, 0xF6, 0xF3)  # 微暖白
INK       = RGBColor(0x1A, 0x1A, 0x1A)  # 纯黑墨
GRAY      = RGBColor(0x66, 0x66, 0x66)  # 中灰
LIGHTGRAY = RGBColor(0xCC, 0xCC, 0xCC)  # 浅灰
ACCENT    = RGBColor(0xC0, 0x39, 0x2B)  # 沉稳红


# ============================================================
#  LLM 两阶段生成
# ============================================================

OUTLINE_PROMPT = """\
你是一个专业的 PPT 内容策划师。根据主题生成一份 PPT 大纲。

要求：
- 8-10 页内容页（不含封面和结尾）
- 每页有明确的 title
- 每页有 3-5 个要点 bullets，每条 15-25 字，精炼有信息量
- 内容由浅入深、逻辑递进
- 第一页 type: "cover"，有 title 和 subtitle
- 最后一页 type: "ending"

严格返回以下 JSON（不要加 markdown 代码块标记）：
{{
  "slides": [
    {{"type": "cover", "title": "...", "subtitle": "..."}},
    {{"type": "content", "title": "...", "bullets": ["...", "...", "..."]}},
    {{"type": "ending", "title": "谢谢", "subtitle": "..."}}
  ]
}}

主题：{query}
语言：{lang}
"""

EXPAND_PROMPT = """\
你是 PPT 内容专家。请为以下幻灯片扩写更详细的内容。

当前页面标题：{title}
当前要点：{bullets}

请在保留原有要点核心意思的基础上，将每条要点扩写为更详细的一句话描述（25-40字），并补充 1-2 条新的要点。
如果该页面适合，可以在适当位置标注可以插入图表/示意图的建议（用 [图表: xxx] 标记）。

严格返回 JSON 数组，每个元素是一个字符串：
["要点1详细描述", "要点2详细描述", "新补充要点", "[图表: xxx]"]

只返回 JSON 数组，不要加其他文字。"""


def _call_llm(messages: list[dict], temperature: float = 0.7) -> str:
    """同步调用星火 LLM"""
    import httpx
    from config import settings

    password = f"{settings.spark_api_key}:{settings.spark_api_secret}"
    payload = {
        "model": spark_service.MODEL,
        "messages": messages,
        "temperature": temperature,
        "max_tokens": 4096,
    }
    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {password}",
    }
    with httpx.Client(timeout=120.0) as client:
        resp = client.post(spark_service._base_url, headers=headers, json=payload)
        resp.raise_for_status()
        data = resp.json()
        return data["choices"][0]["message"]["content"]


def _parse_json_from_llm(content: str):
    """从 LLM 输出中提取 JSON，容错处理"""
    content = content.strip()

    # 剥离 markdown 代码块
    m = re.search(r'```(?:json)?\s*\n([\s\S]*?)```', content)
    if m:
        content = m.group(1).strip()
    elif content.startswith("```"):
        lines = content.split("\n")
        lines = [l for l in lines if not l.strip().startswith("```")]
        content = "\n".join(lines).strip()

    # 找 JSON 边界
    brace_start = content.find('{')
    bracket_start = content.find('[')
    start = -1
    if brace_start != -1 and bracket_start != -1:
        start = min(brace_start, bracket_start)
    elif brace_start != -1:
        start = brace_start
    elif bracket_start != -1:
        start = bracket_start

    if start != -1:
        # 找对应的结束符
        end = max(content.rfind('}'), content.rfind(']'))
        if end > start:
            content = content[start:end + 1]

    # 尝试解析
    try:
        return json.loads(content)
    except json.JSONDecodeError:
        pass

    # 修复单引号
    fixed = content.replace("'", '"')
    fixed = re.sub(r',\s*([}\]])', r'\1', fixed)
    try:
        return json.loads(fixed)
    except json.JSONDecodeError:
        pass

    import ast
    try:
        result = ast.literal_eval(content)
        if isinstance(result, (dict, list)):
            return result
    except (ValueError, SyntaxError):
        pass

    return None


def _escape_format(s: str) -> str:
    """转义 str.format 的占位符花括号"""
    return s.replace("{", "{{").replace("}", "}}")


def _generate_outline(query: str, language: str) -> list[dict]:
    """阶段一：LLM 生成大纲"""
    lang_label = "中文" if language == "cn" else "English"
    prompt = OUTLINE_PROMPT.format(query=_escape_format(query), lang=lang_label)

    raw = _call_llm([{"role": "user", "content": prompt}])
    logger.info("大纲 LLM 原始输出前 300 字: %s", raw[:300])

    result = _parse_json_from_llm(raw)
    if not result or not isinstance(result, dict):
        raise RuntimeError(f"大纲 JSON 解析失败，原始输出:\n{raw[:500]}")

    slides = result.get("slides", [])
    if not slides:
        raise RuntimeError("大纲中没有 slides")

    logger.info("大纲生成成功，共 %d 页", len(slides))
    return slides


def _expand_slide(title: str, bullets: list[str]) -> list[str]:
    """阶段二：LLM 逐页扩写"""
    prompt = EXPAND_PROMPT.format(
        title=_escape_format(title),
        bullets=_escape_format(json.dumps(bullets, ensure_ascii=False)),
    )

    try:
        raw = _call_llm([{"role": "user", "content": prompt}], temperature=0.5)
        result = _parse_json_from_llm(raw)
        if isinstance(result, list) and all(isinstance(x, str) for x in result):
            logger.info("页面「%s」扩写成功，%d 条要点", title, len(result))
            return result
    except Exception as e:
        logger.warning("页面「%s」扩写失败，使用原始要点: %s", title, e)

    return bullets


# ============================================================
#  python-pptx 排版
# ============================================================

def _set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_rect(slide, left, top, width, height, color: RGBColor):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _set_font(run, size: int, bold=False, color=INK, name=FONT_BODY):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name
    # 中文字体回退
    from pptx.oxml.ns import qn
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        from lxml import etree
        ea = etree.SubElement(rPr, qn('a:ea'))
    ea.set('typeface', name)


def _add_cover_slide(prs: Presentation, data: dict):
    """封面：纯白底 + 大标题左对齐 + 红色竖条 + 副标题"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)

    # 左侧红色竖条 — 视觉锚点
    _add_rect(slide, Inches(0.9), Inches(1.8), Pt(5), Inches(2.2), ACCENT)

    # 主标题
    title = data.get("title", "")
    txBox = slide.shapes.add_textbox(Inches(1.3), Inches(1.8), Inches(7.4), Inches(2.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    _set_font(run, 40, bold=True, color=INK)

    # 副标题
    subtitle = data.get("subtitle", "")
    if subtitle:
        txSub = slide.shapes.add_textbox(Inches(1.3), Inches(4.2), Inches(7.4), Inches(0.8))
        p2 = txSub.text_frame.paragraphs[0]
        run2 = p2.add_run()
        run2.text = subtitle
        _set_font(run2, 16, color=GRAY)

    # 底部细线
    _add_rect(slide, Inches(1.3), Inches(6.4), Inches(1.5), Pt(1.5), LIGHTGRAY)


def _add_content_slide(prs: Presentation, data: dict, page_num: int, total: int):
    """内容页：白底 + 标题区 + 分隔线 + 要点列表 + 页码"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)

    # 页码 — 右上角
    txPage = slide.shapes.add_textbox(Inches(8.5), Inches(0.4), Inches(1.0), Inches(0.4))
    pPage = txPage.text_frame.paragraphs[0]
    pPage.alignment = PP_ALIGN.RIGHT
    run_page = pPage.add_run()
    run_page.text = f"{page_num:02d}/{total:02d}"
    _set_font(run_page, 9, color=LIGHTGRAY)

    # 标题
    title = data.get("title", "")
    txTitle = slide.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(8.0), Inches(0.8))
    tf = txTitle.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run_t = p.add_run()
    run_t.text = title
    _set_font(run_t, 26, bold=True, color=INK)

    # 标题下方红色细线
    _add_rect(slide, Inches(1.0), Inches(1.4), Inches(0.8), Pt(2.5), ACCENT)

    # 要点列表
    bullets = data.get("bullets", [])
    if not bullets:
        return

    txBody = slide.shapes.add_textbox(Inches(1.0), Inches(1.9), Inches(8.0), Inches(5.0))
    tfBody = txBody.text_frame
    tfBody.word_wrap = True

    for i, bullet in enumerate(bullets):
        p = tfBody.paragraphs[0] if i == 0 else tfBody.add_paragraph()
        p.space_after = Pt(10)
        p.space_before = Pt(2)

        # 检查是否是图表建议
        if bullet.startswith("[图表:") and bullet.endswith("]"):
            chart_text = bullet[3:-1]  # 去掉 [图表: 和 ]
            run_mark = p.add_run()
            run_mark.text = f"📊  {chart_text}"
            _set_font(run_mark, 11, color=GRAY)
            continue

        # 序号
        run_num = p.add_run()
        run_num.text = f"  {i + 1}.  "
        _set_font(run_num, 13, bold=True, color=ACCENT)

        # 要点文字
        run_text = p.add_run()
        run_text.text = bullet
        _set_font(run_text, 14, color=INK)


def _add_ending_slide(prs: Presentation, data: dict):
    """结尾页：居中致谢"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)

    title = data.get("title", "谢谢")
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    _set_font(run, 44, bold=True, color=INK)

    subtitle = data.get("subtitle", "")
    if subtitle:
        txSub = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(0.6))
        p2 = txSub.text_frame.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = subtitle
        _set_font(run2, 14, color=GRAY)

    # 居中红色短线
    _add_rect(slide, Inches(4.5), Inches(5.0), Inches(1.0), Pt(2.5), ACCENT)


# ============================================================
#  主生成函数
# ============================================================

def generate_local_ppt(query: str, language: str = "cn") -> str:
    """本地生成 PPT，返回文件路径"""
    logger.info("本地 PPT 生成开始，主题: %s", query)

    # 阶段一：生成大纲
    slides_data = _generate_outline(query, language)

    # 阶段二：逐页扩写内容
    content_slides = [s for s in slides_data if s.get("type") == "content"]
    for i, sd in enumerate(content_slides):
        _local_tasks_ref = None  # 用于更新进度
        expanded = _expand_slide(sd["title"], sd.get("bullets", []))
        sd["bullets"] = expanded
        # 更新进度：30% + (i/total)*50%
        progress = 30 + int((i + 1) / len(content_slides) * 50)
        if hasattr(generate_local_ppt, '_current_task_id'):
            tid = generate_local_ppt._current_task_id
            if tid in _local_tasks:
                _local_tasks[tid]["progress"] = progress

    # 组装 PPT
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    total_content = len(content_slides)
    page_num = 1
    for sd in slides_data:
        t = sd.get("type", "content")
        if t == "cover":
            _add_cover_slide(prs, sd)
        elif t == "ending":
            _add_ending_slide(prs, sd)
        else:
            _add_content_slide(prs, sd, page_num, total_content)
            page_num += 1

    # 保存
    task_id = uuid.uuid4().hex[:12]
    tmp_dir = os.path.join(tempfile.gettempdir(), "smartlearn_ppt")
    os.makedirs(tmp_dir, exist_ok=True)
    filepath = os.path.join(tmp_dir, f"ppt_{task_id}.pptx")
    prs.save(filepath)

    logger.info("本地 PPT 已保存: %s", filepath)
    return filepath


# ============================================================
#  异步任务管理
# ============================================================
_local_tasks: dict[str, dict] = {}


def create_local_task(query: str, language: str) -> str:
    """创建本地 PPT 生成任务，返回 task_id"""
    import threading

    task_id = uuid.uuid4().hex[:12]
    _local_tasks[task_id] = {"status": "generating", "progress": 0, "file_path": "", "error": ""}

    def _run():
        try:
            _local_tasks[task_id]["progress"] = 10
            generate_local_ppt._current_task_id = task_id
            filepath = generate_local_ppt(query, language)
            _local_tasks[task_id].update(status="done", progress=100, file_path=filepath)
        except Exception as e:
            logger.error("本地 PPT 生成失败: %s", e, exc_info=True)
            _local_tasks[task_id].update(status="error", error=str(e))

    t = threading.Thread(target=_run, daemon=True)
    t.start()
    return task_id


def get_local_task(task_id: str) -> dict | None:
    return _local_tasks.get(task_id)
